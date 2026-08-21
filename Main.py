import io
import os
import zipfile
import logging

from fastapi import FastAPI, Request, HTTPException
from fastapi.responses import JSONResponse
from pptx import Presentation


# ============================================================
# APP
# ============================================================

app = FastAPI(
    title="Tata RTSM PPT Extractor",
    version="4.0"
)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s"
)

logger = logging.getLogger("rtms-ppt-extractor")


# ============================================================
# SETTINGS
# ============================================================

MAX_FILE_SIZE = 50 * 1024 * 1024

PPTX_SIGNATURE = b"PK\x03\x04"

PPT_SIGNATURE = bytes.fromhex(
    "D0CF11E0A1B11AE1"
)


# ============================================================
# HOME
# ============================================================

@app.get("/")
async def home():

    return {
        "status": "live",
        "service": "Tata RTSM PPT Extractor",
        "version": "4.0",
        "accepted_files": [
            ".pptx",
            ".ppt"
        ],
        "endpoint": "/extract"
    }


# ============================================================
# HEALTH
# ============================================================

@app.get("/health")
async def health():

    return {
        "status": "healthy"
    }


# ============================================================
# FIND PPTX ZIP SIGNATURE
# ============================================================

def find_pptx_signature(data: bytes):

    if not data:
        return -1

    # Normal PPTX
    if data.startswith(PPTX_SIGNATURE):
        return 0

    # Search first 1 MB only
    # to avoid accidentally treating random binary
    # content as a PPTX.
    position = data[:1024 * 1024].find(
        PPTX_SIGNATURE
    )

    return position


# ============================================================
# NORMALIZE PPTX
# ============================================================

def normalize_pptx(data: bytes):

    position = find_pptx_signature(data)

    if position == -1:

        raise ValueError(
            "PowerPoint ZIP signature PK0304 "
            "was not found."
        )

    if position > 0:

        logger.warning(
            "Found %d byte(s) before PPTX ZIP signature. "
            "Removing them.",
            position
        )

        data = data[position:]

    return data


# ============================================================
# DETECT FILE TYPE
# ============================================================

def detect_file_type(data: bytes):

    if not data:
        return "unknown"

    # PPTX
    if find_pptx_signature(data) >= 0:
        return "pptx"

    # Old binary PPT
    if data.startswith(PPT_SIGNATURE):
        return "ppt"

    return "unknown"


# ============================================================
# VALIDATE PPTX
# ============================================================

def validate_pptx(data: bytes):

    logger.info(
        "Starting PPTX validation..."
    )

    normalized = normalize_pptx(data)

    logger.info(
        "Normalized PPTX size: %d bytes",
        len(normalized)
    )

    logger.info(
        "Normalized first 16 bytes: %s",
        normalized[:16].hex()
    )

    # --------------------------------------------------------
    # Check ZIP
    # --------------------------------------------------------

    try:

        is_zip = zipfile.is_zipfile(
            io.BytesIO(normalized)
        )

    except Exception as exc:

        logger.exception(
            "ZIP detection exception: %s",
            exc
        )

        raise ValueError(
            f"ZIP detection failed: {exc}"
        )

    if not is_zip:

        raise ValueError(
            "The supplied file contains a PPTX signature "
            "but is not recognised as a valid ZIP package."
        )

    # --------------------------------------------------------
    # Open ZIP
    # --------------------------------------------------------

    try:

        z = zipfile.ZipFile(
            io.BytesIO(normalized),
            mode="r",
            allowZip64=True
        )

        names = z.namelist()

        logger.info(
            "ZIP opened successfully."
        )

        logger.info(
            "ZIP contains %d entries.",
            len(names)
        )

    except zipfile.BadZipFile as exc:

        logger.error(
            "BAD ZIP ERROR: %s",
            str(exc)
        )

        raise ValueError(
            f"PPTX ZIP package is corrupted or incomplete: {exc}"
        )

    except Exception as exc:

        logger.exception(
            "ZIP opening exception."
        )

        raise ValueError(
            f"Unable to open PPTX package: {exc}"
        )

    # --------------------------------------------------------
    # Verify PowerPoint package
    # --------------------------------------------------------

    if "[Content_Types].xml" not in names:

        z.close()

        raise ValueError(
            "ZIP package does not contain "
            "[Content_Types].xml. "
            "This is not a valid PPTX package."
        )

    if "ppt/presentation.xml" not in names:

        z.close()

        raise ValueError(
            "ZIP package does not contain "
            "ppt/presentation.xml. "
            "This is not a valid PowerPoint file."
        )

    z.close()

    logger.info(
        "PPTX package structure validated successfully."
    )

    return normalized


# ============================================================
# EXTRACT TEXT
# ============================================================

def extract_pptx_text(data: bytes):

    logger.info(
        "Opening PowerPoint with python-pptx..."
    )

    try:

        presentation = Presentation(
            io.BytesIO(data)
        )

    except Exception as exc:

        logger.exception(
            "python-pptx failed."
        )

        raise ValueError(
            f"PowerPoint could not be opened: {exc}"
        )

    slides_output = []

    all_text = []

    slide_count = len(
        presentation.slides
    )

    logger.info(
        "PowerPoint contains %d slides.",
        slide_count
    )

    # --------------------------------------------------------
    # SLIDES
    # --------------------------------------------------------

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        slide_parts = []

        for shape in slide.shapes:

            # ----------------------------------------------
            # Text
            # ----------------------------------------------

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:

                    slide_parts.append(
                        text
                    )

            # ----------------------------------------------
            # Tables
            # ----------------------------------------------

            if getattr(
                shape,
                "has_table",
                False
            ):

                table = shape.table

                for row in table.rows:

                    values = []

                    for cell in row.cells:

                        value = cell.text.strip()

                        if value:

                            values.append(
                                value
                            )

                    if values:

                        slide_parts.append(
                            " | ".join(values)
                        )

        slide_text = "\n".join(
            slide_parts
        )

        slides_output.append(
            {
                "slide_number": slide_number,
                "text": slide_text
            }
        )

        if slide_text:

            all_text.append(
                f"SLIDE {slide_number}\n"
                f"{slide_text}"
            )

    return {
        "slide_count": slide_count,
        "slides": slides_output,
        "text": "\n\n".join(all_text)
    }


# ============================================================
# READ REQUEST
# ============================================================

async def get_request_bytes(
    request: Request
):

    body = await request.body()

    if not body:

        raise HTTPException(
            status_code=400,
            detail="No file content received."
        )

    if len(body) > MAX_FILE_SIZE:

        raise HTTPException(
            status_code=413,
            detail=(
                "File is larger than "
                f"{MAX_FILE_SIZE // (1024 * 1024)} MB."
            )
        )

    return body


# ============================================================
# MAIN EXTRACT API
# ============================================================

@app.post("/extract")
async def extract(request: Request):

    logger.info("")
    logger.info("=" * 70)
    logger.info("RTSM PPT EXTRACT REQUEST STARTED")
    logger.info("=" * 70)

    try:

        # ----------------------------------------------------
        # READ BODY
        # ----------------------------------------------------

        body = await get_request_bytes(
            request
        )

        content_type = request.headers.get(
            "content-type",
            ""
        )

        content_length = request.headers.get(
            "content-length",
            ""
        )

        logger.info(
            "Content-Type: %s",
            content_type
        )

        logger.info(
            "Content-Length: %s",
            content_length
        )

        logger.info(
            "Received file size: %d bytes",
            len(body)
        )

        logger.info(
            "First 64 bytes: %s",
            body[:64].hex()
        )

        # ----------------------------------------------------
        # FILE TYPE
        # ----------------------------------------------------

        file_type = detect_file_type(
            body
        )

        logger.info(
            "Detected file type: %s",
            file_type
        )

        # ----------------------------------------------------
        # ONLY POWERPOINT
        # ----------------------------------------------------

        if file_type == "unknown":

            logger.error(
                "REJECTED: File is not PowerPoint."
            )

            raise HTTPException(
                status_code=400,
                detail=(
                    "Only .pptx PowerPoint files are accepted."
                )
            )

        # ----------------------------------------------------
        # OLD PPT
        # ----------------------------------------------------

        if file_type == "ppt":

            logger.error(
                "Old .ppt file received."
            )

            return JSONResponse(
                status_code=400,
                content={
                    "success": False,
                    "error": (
                        "Old .ppt format is not supported. "
                        "Please convert the file to .pptx."
                    )
                }
            )

        # ----------------------------------------------------
        # VALIDATE
        # ----------------------------------------------------

        pptx_data = validate_pptx(
            body
        )

        # ----------------------------------------------------
        # EXTRACT
        # ----------------------------------------------------

        result = extract_pptx_text(
            pptx_data
        )

        # ----------------------------------------------------
        # SUCCESS
        # ----------------------------------------------------

        logger.info(
            "PPTX EXTRACTION SUCCESSFUL."
        )

        logger.info(
            "Slides extracted: %d",
            result["slide_count"]
        )

        logger.info("=" * 70)
        logger.info("RTSM PPT EXTRACT REQUEST COMPLETED")
        logger.info("=" * 70)

        return {
            "success": True,
            "file_type": "pptx",
            "file_size": len(body),
            "slide_count": result["slide_count"],
            "slides": result["slides"],
            "text": result["text"]
        }

    except HTTPException:

        raise

    except ValueError as exc:

        logger.error(
            "VALIDATION ERROR: %s",
            str(exc)
        )

        raise HTTPException(
            status_code=400,
            detail=str(exc)
        )

    except Exception as exc:

        logger.exception(
            "UNEXPECTED ERROR"
        )

        raise HTTPException(
            status_code=500,
            detail=(
                "PPT extraction failed: "
                + str(exc)
            )
        )


# ============================================================
# START SERVER
# ============================================================

if __name__ == "__main__":

    import uvicorn

    port = int(
        os.environ.get(
            "PORT",
            "10000"
        )
    )

    uvicorn.run(
        "main:app",
        host="0.0.0.0",
        port=port
    )
